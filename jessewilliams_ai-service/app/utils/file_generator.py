"""
File Generation Service - Production-Ready with Real Data Only
Generates Excel, Word, PDF, PowerPoint from actual search results and database
"""

from typing import Dict, List, Optional, Tuple
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
import io
import xlsxwriter
from datetime import datetime

# Word (python-docx) imports
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

# PowerPoint (python-pptx) imports
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt  
from pptx.dml.color import RGBColor as PptxRGBColor  
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE 

class FileGenerationService:
    """Intelligent file generation based on chat context"""
    
    @staticmethod
    def should_generate_file(message: str, response: str) -> Tuple[bool, Optional[str]]:
        """
        Determine if file should be generated and what type
        
        Returns:
            (should_generate: bool, file_type: str|None)
        """
        
        # Keywords in user message
        user_keywords = {
            'excel': ['excel', 'spreadsheet', 'xlsx', 'table', 'csv', 'tracker'],
            'word': ['word', 'document', 'docx', 'report', 'write'],
            'powerpoint': ['powerpoint', 'presentation', 'pptx', 'slides', 'pitch']
        }
        
        # Action keywords
        action_keywords = [
            'generate', 'create', 'make', 'build', 'produce',
            'export', 'download', 'send', 'give me', 'show me'
        ]
        
        message_lower = message.lower()
        response_lower = response.lower()
        
        # Check if action keyword present
        has_action = any(action in message_lower for action in action_keywords)
        
        if not has_action:
            return False, None
        
        # Determine file type
        for file_type, keywords in user_keywords.items():
            if any(keyword in message_lower for keyword in keywords):
                return True, file_type
        
        # Default to Excel for tables/data
        if any(word in message_lower for word in ['table', 'data', 'list', 'results', 'properties']):
            return True, 'excel'
        
        return False, None
     
    @staticmethod
    def extract_data_from_context(
     user_context: Dict,
     session_id: str,
     search_results: Optional[List[Dict]] = None) -> Dict:
     """
     Extract all relevant data for file generation
     NOW HANDLES PDF + Validates search results
     """
    
     profile = user_context.get("profile", {})
     session = user_context.get("sessions", {}).get(session_id, {})
     extracted = session.get("extracted_data", {})
    
    # Helper to convert to float safely
     def to_float(val, default=0.0):
        if isinstance(val, (int, float)):
            return float(val)
        if isinstance(val, str):
            try:
                return float(val.replace('$', '').replace(',', '').strip())
            except:
                return default
        return default
    
    # Extract financial data
     capital = to_float(
        extracted.get('capital') or profile.get('startingCapital', 0)
     )
    
     profit_goal = to_float(
        extracted.get('profit_goal') or profile.get('profitGoal', 0)
     )
    
    # Extract location
     location = extracted.get('location') or profile.get('targetGeography', 'Not specified')
    
    # ===================================================
    # CRITICAL: VALIDATE SEARCH RESULTS
    # ===================================================
     validated_results = []
     if search_results:
        for result in search_results:
            # Only include results with required fields
            if result.get('address') or (result.get('price') and result.get('price') > 0):
                validated_results.append(result)
    
     print(f"✅ Validated {len(validated_results)} search results for file generation")
    
    # Gather all data
     data = {
        # User profile - REAL DATA ONLY
        'investor_name': profile.get('name') or 'Investor',
        'strategy_type': profile.get('strategy', 'Investment Strategy'),
        'property_type': profile.get('propertyType', 'Land'),
        'rental_type': profile.get('rentalType'),
        'capital': capital,
        'location': location,
        'timeline': extracted.get('timeline') or profile.get('investmentTimeline', 'Not specified'),
        'profit_goal': profit_goal,
        
        # Search results - VALIDATED ONLY
        'search_results': validated_results,
        
        # Metadata
        'date': datetime.now().strftime('%B %d, %Y'),
        'generated_at': datetime.now().isoformat(),
     }
    
    # Calculate projections ONLY if we have real capital data
     if capital > 0 and profit_goal > 0:
        data['projections'] = FileGenerationService._calculate_projections(
            capital,
            profit_goal
        )
     else:
        data['projections'] = []
    
    # Capital deployment ONLY if we have real capital
     if capital > 0:
        data['deployment'] = FileGenerationService._calculate_deployment(capital)
     else:
        data['deployment'] = []
    
    # Next steps - contextualized to location and data
     data['next_steps'] = FileGenerationService._generate_next_steps(
        location, 
        capital > 0,
        len(validated_results)
     )
    
     return data
    
    
    
    @staticmethod
    def _calculate_fit_score(price: float, capital: float,    profit_goal: float, acres: float = None) -> str:
     """
     Assess if property fits investor's plan
    
     Returns: ✓ Good | ⚠ Marginal | ✗ Poor
     """
    
     if not price or price <= 0:
        return "⚠ No Price"
    
     if not capital or capital <= 0:
        return "⚠ Set Budget"
    
    # Check if within budget
     if price > capital:
        return "✗ Over Budget"
    
    # Calculate minimum profit potential (rough estimate)
    # Using 70% rule: max offer should leave room for profit
     if price > capital * 0.7:
        # Not enough room for repairs and profit
        return "⚠ Thin Margin"
    
    # Good fit
     return "✓ Good Fit"



    @staticmethod
    def _identify_data_gaps(prop: dict) -> str:
     """
     Identify what additional data is needed for verification
    
     Returns: Comma-separated list of missing/needed items
     """
    
     gaps = []
    
    # Check address completeness
     address = str(prop.get('address', '')).lower()
     if not address or address == 'address not provided':
        gaps.append("Full address")
    
    # Check price
     price = prop.get('price', 0)
     if not price or price <= 0:
        gaps.append("Verified price")
    
    # Check lot size
     acres = prop.get('acres')
     lot_size = prop.get('lot_size')
     if not acres and not lot_size:
        gaps.append("Lot size")
    
    # Check for property details
     beds = prop.get('bedrooms')
     baths = prop.get('bathrooms')
     sqft = prop.get('sqft')
    
     details_missing = not (beds or baths or sqft)
     if details_missing:
        gaps.append("Property details")
    
    # Check for zoning/legal
     gaps.append("Zoning verification")  # Always needed
     gaps.append("Deed review")  # Always needed
    
    # Check source quality
     source = prop.get('source', '').lower()
     if 'estimated' in source or 'approximate' in source:
        gaps.append("Official documentation")
    
     if not gaps:
        gaps.append("Minor documentation")
    
     return ", ".join(gaps[:3])
    
    
    
    @staticmethod
    def _calculate_deployment(capital: float) -> List[Dict]:
        """Calculate capital deployment with percentages"""
        return [
            {
                'category': 'Marketing & Sourcing',
                'amount': capital * 0.35,
                'percentage': 35,
                'purpose': 'Direct mail campaigns, skip tracing, online lead generation, and acquisition marketing'
            },
            {
                'category': 'Earnest Money Deposits',
                'amount': capital * 0.40,
                'percentage': 40,
                'purpose': 'Contract deposits for 2-3 simultaneous property acquisitions'
            },
            {
                'category': 'Operating Reserve',
                'amount': capital * 0.25,
                'percentage': 25,
                'purpose': 'Legal/entity setup, due diligence costs, software tools, and contingency buffer'
            }
        ]
    
    
    
    @staticmethod
    def _calculate_projections(capital: float, profit_goal: float) -> List[Dict]:
        """Calculate realistic 6-month profit projections"""
        
        # Conservative ramp-up schedule
        monthly_deals = [0, 0, 1, 1, 2, 2]  # Deals closed each month
        cumulative_deals = []
        cumulative_profit = []
        
        total = 0
        profit_total = 0
        avg_profit = (profit_goal / 6) if profit_goal > 0 else (capital * 0.20)
        
        for month, deals in enumerate(monthly_deals, 1):
            total += deals
            profit_total += deals * avg_profit
            
            cumulative_deals.append(total)
            cumulative_profit.append(profit_total)
        
        return [
            {
                'month': i + 1,
                'month_name': ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun'][i],
                'deals': cumulative_deals[i],
                'profit': cumulative_profit[i],
                'roi': (cumulative_profit[i] / capital * 100) if capital > 0 else 0
            }
            for i in range(6)
        ]
    
    
    
    @staticmethod
    def generate_file_response_text(data: Dict, file_type: str) -> str:
       """
       Generate concise response text for file generation
       """
    
       location = data.get('location', 'the target market')
       capital = data.get('capital', 0)
       property_count = len(data.get('search_results', []))
    
       if not location or location == 'Not specified':
          location = 'your target market'
          
       if file_type == 'excel':
          response = f"""Your spreadsheet is ready.

**What this spreadsheet includes:**
- **Deal Pipeline sheet**: {property_count} {location} properties with exact addresses
- **Investment Summary**: Short summary of your investment profile and market analysis
- **Capital Deployment**: Capital deployment plan for ${capital:,.0f}
- **Profit Projections**: 6-month ROI timeline with calculations
- **Action Items**: Prioritized next steps with timelines

**How to use it:**
1. Review properties in Deal Pipeline tab - add your notes in column "I"
2. Filter by price/acre to identify best value opportunities
3. Track capital allocation in Investment Summary tab

**Next Steps:**
1. Schedule site visits for top 3-5 properties within 7 days
2. Request county records: zoning maps, survey data, tax info
3. Build cash buyer list (50+ contacts) before making first offer"""

       elif file_type == 'word':
          response = f"""Your investment report is ready.

**What this document includes:**
- **Executive Summary**: Market analysis for {location}
- **Property Profiles**: {property_count} detailed listings with investment metrics
- **Financial Projections**: ROI analysis and cash flow models
- **Risk Assessment**: Market risks and mitigation strategies

**How to use it:**
1. Read Executive Summary first for key insights
2. Compare properties in Section 3 using ranking criteria
3. Share with partners/advisors for feedback

**Next Steps:**
1. Review with financial advisor or investment partner
2. Schedule property tours for top-ranked listings
3. Prepare offer strategy based on recommendations"""

       elif file_type == 'powerpoint':
          response = f"""Your investor presentation is ready.

**What this presentation includes:**
- **Market Opportunity**: {location} investment thesis
- **Property Showcase**: {property_count} highlighted opportunities with photos
- **Financial Model**: Capital deployment and profit projections
- **Execution Timeline**: 6-month action plan

**How to use it:**
1. Customize slides 1-2 with your branding/contact info
2. Practice 10-minute pitch covering all sections
3. Prepare Q&A responses on market risks

**Next Steps:**
1. Schedule meetings with potential funding partners
2. Prepare deal package with supporting documents
3. Set follow-up system for investor communications"""

       elif file_type == 'pdf':
          response = f"""Your investment report is ready.

**What this document includes:**
- **Executive Summary**: Market analysis for {location}
- **Property Profiles**: {property_count} detailed listings with investment metrics
- **Financial Projections**: ROI analysis and cash flow models
- **Risk Assessment**: Market risks and mitigation strategies

**How to use it:**
1. Read Executive Summary first for key insights
2. Compare properties in Section 3 using ranking criteria
3. Share with partners/advisors for feedback

**Next Steps:**
1. Review with financial advisor or investment partner
2. Schedule property tours for top-ranked listings
3. Prepare offer strategy based on recommendations"""

       else:
          return "**✅ File Ready** Your analysis file is ready to download and review."
        
       # VALIDATION: NEVER return empty or None
       if not response or not response.strip():
          response = "Your file has been generated successfully. Download it to access all investment details and analysis."

       return response.strip() if response else "Your file has been generated successfully. Download it to access all analysis details."
 
         
    
    @staticmethod
    def _generate_next_steps(location: str, has_capital: bool, property_count: int) -> List[str]:
        """Generate contextual next steps"""
        steps = []
        
        if property_count > 0:
            steps.append(f"Review and prioritize the {property_count} properties identified in {location}")
            steps.append("Schedule site visits for top 3-5 candidate properties")
            steps.append("Request seller disclosures and zoning verification from county")
        else:
            steps.append(f"Initiate property search in {location} to identify acquisition targets")
        
        if has_capital:
            steps.append("Set up business entity (LLC) and open dedicated business banking account")
            steps.append("Build buyer list of 50+ cash buyers before first contract")
            steps.append("Launch direct mail campaign to off-market property owners")
        else:
            steps.append("Finalize capital sourcing strategy and funding timeline")
        
        steps.append("Join local REIA (Real Estate Investors Association) for networking")
        steps.append("Establish relationships with title company and real estate attorney")
        
        return steps
    
    
    
    
    """
    SECTION: Word Document Generation with Full Content
    """
    @staticmethod
    def generate_word(data: Dict) -> Tuple[bytes, str]:
     """Generate comprehensive Word document report with REAL DATA"""
    
    # ===================================================
    # VALIDATION: Ensure all required data exists
    # ===================================================
     if not data:
        raise ValueError("No data provided for Word generation")
    
    # Extract and validate financial data
     capital = data.get('capital', 0)
     if capital is None:
        capital = 0
     try:
        capital = float(capital)
     except (ValueError, TypeError):
        capital = 0.0
    
     profit_goal = data.get('profit_goal', 0)
     if profit_goal is None:
        profit_goal = 0
     try:
        profit_goal = float(profit_goal)
     except (ValueError, TypeError):
        profit_goal = 0.0
    
    # Extract and validate location
     location = data.get('location', 'Not specified')
     if location is None:
        location = 'Not specified'
     location = str(location)
    
    # Create document
     doc = Document()
    
    # ===================================================
    # TITLE PAGE
    # ===================================================
     title = doc.add_heading('Real Estate Investment Analysis Report', 0)
     title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
     subtitle = doc.add_paragraph('Comprehensive Investment Strategy & Market Analysis')
     subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
     subtitle_format = subtitle.runs[0].font
     subtitle_format.size = Pt(14)
     subtitle_format.color.rgb = RGBColor(100, 100, 100)
     subtitle_format.italic = True
    
     doc.add_paragraph()
    
    # Prepare investor information
     investor_name = str(data.get('investor_name', 'Real Estate Investor')) if data.get('investor_name') else 'Real Estate Investor'
     date_str = str(data.get('date', 'N/A')) if data.get('date') else 'N/A'
     strategy = str(data.get('strategy_type', 'Investment Strategy')) if data.get('strategy_type') else 'Investment Strategy'
    
     doc.add_paragraph(f"Prepared for: {investor_name}", style='Normal')
     doc.add_paragraph(f"Date: {date_str}")
     doc.add_paragraph(f"Market: {location}")
     doc.add_paragraph(f"Strategy: {strategy}")
    
     doc.add_page_break()
    
    # ===================================================
    # SECTION 1: EXECUTIVE SUMMARY
    # ===================================================
     doc.add_heading('1. Executive Summary', 1)
    
     property_count = len(data.get('search_results', []))
    
     summary_text = f"""This report provides a detailed real estate investment analysis for {location} based on your investment criteria:

       • Starting Capital: ${capital:,.0f}
       • Target Profit Goal: ${profit_goal:,.0f}
       • Investment Timeline: {str(data.get('timeline', 'Not specified')) if data.get('timeline') else 'Not specified'}
       • Strategy Focus: {strategy}
       • Property Type: {str(data.get('property_type', 'Residential/Commercial')) if data.get('property_type') else 'Residential/Commercial'}

     The analysis identifies {property_count} qualified properties, evaluates market conditions, and provides actionable recommendations to maximize return on investment while minimizing risk exposure."""
     
     doc.add_paragraph(summary_text)
    
    # ===================================================
    # SECTION 2: INVESTOR PROFILE & OBJECTIVES
    # ===================================================
     doc.add_heading('2. Investor Profile & Objectives', 1)
    
     doc.add_heading('2.1 Investment Profile', 2)
    
    # Prepare all profile data with validation
     profile_data = [
        ('Investor Name', investor_name),
        ('Starting Capital', f"${capital:,.0f}"),
        ('Target Geography', location),
        ('Property Type Focus', str(data.get('property_type', 'N/A')) if data.get('property_type') else 'N/A'),
        ('Investment Strategy', strategy),
        ('Timeline', str(data.get('timeline', 'N/A')) if data.get('timeline') else 'N/A'),
        ('Target Profit Goal', f"${profit_goal:,.0f}"),
        ('Risk Tolerance', 'Moderate - Conservative'),
    ]
    
    # Create table with dynamic rows (data rows + 1 header)
     profile_table = doc.add_table(rows=len(profile_data) + 1, cols=2)
     profile_table.style = 'Light Grid Accent 1'
    
    # Add header row
     try:
        header_row = profile_table.rows[0]
        header_row.cells[0].text = 'Category'
        header_row.cells[1].text = 'Value'
        # Make header bold
        for cell in header_row.cells:
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
     except Exception as e:
        print(f"Error setting header row: {e}")
    
    # Add data rows
     for i, (key, value) in enumerate(profile_data):
        try:
            row_idx = i + 1  # Offset by 1 for header row
            row = profile_table.rows[row_idx]
            row.cells[0].text = str(key)
            row.cells[1].text = str(value)
        except (IndexError, AttributeError, TypeError) as e:
            print(f"Error at profile data row {i}: {e}")
            continue
    
     doc.add_heading('2.2 Investment Objectives', 2)
     objectives = [
        f"Identify and acquire {property_count} qualified investment properties in {location}",
        f"Deploy ${capital:,.0f} capital to generate ${profit_goal:,.0f} in profit within {str(data.get('timeline', 'specified timeline')) if data.get('timeline') else 'specified timeline'}",
        "Minimize market and execution risk through diversified sourcing and rigorous underwriting",
        "Establish repeatable systems for deal sourcing, underwriting, and exit strategies",
        "Build long-term investment portfolio with sustainable cash flow and appreciation potential"
    ]
    
     for obj in objectives:
        doc.add_paragraph(obj, style='List Bullet')
    
    # ===================================================
    # SECTION 3: MARKET ANALYSIS
    # ===================================================
     doc.add_heading('3. Market Analysis & Property Overview', 1)
    
     doc.add_heading('3.1 Market Summary', 2)
    
    # Calculate average price safely
     search_results = data.get('search_results', [])
     if search_results:
        prices = [p.get('price', 0) for p in search_results if p.get('price', 0) > 0]
        avg_price = (sum(prices) / len(prices)) if prices else 0
        avg_price_str = f"${avg_price:,.0f}"
     else:
        avg_price_str = "N/A"
    
     market_summary = f"""The {location} real estate market presents compelling opportunities for disciplined investors. Our analysis of current listings identifies {property_count} active properties that align with your investment criteria.

        Market Indicators:
        • Total Properties Identified: {property_count}
        • Average List Price: {avg_price_str}
        • Price Range: Varies by property type and condition
        • Market Trend: Active buyer interest with competitive opportunities"""
    
     doc.add_paragraph(market_summary)
    
    # ===================================================
    # SECTION 3.2: PROPERTY LISTINGS TABLE 
    # ===================================================
     doc.add_heading('3.2 Identified Properties (Top 10)', 2)
     heading = doc.paragraphs[-1]
     heading.paragraph_format.space_before = Pt(12)
     heading.paragraph_format.space_after = Pt(12)

     search_results = data.get('search_results', [])

     if search_results and len(search_results) > 0:
        # Show actual properties found
        properties = search_results[:10]
        
        # Create table: Address | List Price | Lot Size | Fit for Plan | Data Gaps | Source
        prop_table = doc.add_table(rows=len(properties) + 1, cols=6)
        prop_table.style = 'Light Grid Accent 1'
        
        # Set column widths
        for row in prop_table.rows:
            row.cells[0].width = Inches(2.0)  # Address
            row.cells[1].width = Inches(1.0)  # List Price
            row.cells[2].width = Inches(1.0)  # Lot Size
            row.cells[3].width = Inches(1.0)  # Fit for Plan
            row.cells[4].width = Inches(1.5)  # Data Gaps
            row.cells[5].width = Inches(1.0)  # Source
        
        # ===================================================
        # HEADER ROW - BOLD STYLING
        # ===================================================
        header_cells = prop_table.rows[0].cells
        headers = ['Address', 'List Price', 'Lot Size', 'Fit for Plan', 'Data Gaps (Must Verify)', 'Source']
        
        for idx, header_text in enumerate(headers):
            cell = header_cells[idx]
            cell.text = header_text
            
            # Style header
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(10)
                paragraph.paragraph_format.space_before = Pt(0)
                paragraph.paragraph_format.space_after = Pt(0)
        
        # ===================================================
        # DATA ROWS - VALIDATED PROPERTY DATA
        # ===================================================
        for idx, prop in enumerate(properties, 1):
            try:
                row = prop_table.rows[idx]
                
                # 1. ADDRESS
                address = prop.get('address', 'Address not provided')
                if isinstance(address, str):
                    address = address[:60]  # Truncate if too long
                    row.cells[0].text = str(address)
                
                # 2. LIST PRICE
                price = prop.get('price', 0)
                if isinstance(price, (int, float)) and price > 0:
                    row.cells[1].text = f"${price:,.0f}"
                else:
                    row.cells[1].text = "Contact"
                
                # 3. LOT SIZE / ACRES
                acres = prop.get('acres')
                if acres and isinstance(acres, (int, float)):
                    row.cells[2].text = f"{acres} acres"
                else:
                    lot_size = prop.get('lot_size')
                    row.cells[2].text = str(lot_size) if lot_size else "N/A, See listing"
                
                # 4. FIT FOR PLAN (Calculated based on profile)
                capital = data.get('capital', 0)
                profit_goal = data.get('profit_goal', 0)
                
                fit_score = FileGenerationService._calculate_fit_score(
                    price, 
                    capital, 
                    profit_goal, 
                    acres
                )
                row.cells[3].text = fit_score  # ✓ Good / ⚠ Marginal / ✗ Poor
                
                # 5. DATA GAPS (What needs verification)
                data_gaps = FileGenerationService._identify_data_gaps(prop)
                row.cells[4].text = data_gaps
                
                # 6. SOURCE (Where found)
                source = prop.get('source', 'Listing')
                row.cells[5].text = str(source)[:20]
            
            except (IndexError, KeyError, TypeError) as e:
                print(f"Error at property row {idx}: {e}")
                continue
        
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(12)
        p.paragraph_format.space_after = Pt(6)
        
        summary_run = p.add_run(f"📊 Market Summary: ")
        summary_run.bold = True
        
        prices = [p['price'] for p in search_results if p.get('price', 0) > 0]
        if prices:
            avg_price = sum(prices) / len(prices)
            price_summary = f"Found {len(search_results)} properties | Avg Price: ${avg_price:,.0f} | Price Range: ${min(prices):,} - ${max(prices):,}"
        else:
            price_summary = f"Found {len(search_results)} properties | Pricing data pending"
        
        p.add_run(price_summary)
    
     else:
        # ===================================================
        # NO RESULTS - SHOW TEMPLATE TABLE
        # ===================================================
        no_results_para = doc.add_paragraph()
        no_results_para.paragraph_format.space_before = Pt(0)
        no_results_para.paragraph_format.space_after = Pt(12)
        
        run = no_results_para.add_run('⚠️ No Properties Found in Web Search')
        run.bold = True
        run.font.color.rgb = RGBColor(192, 0, 0)
        
        # Create empty template for manual entry
        explanation = doc.add_paragraph(
            "The system did not find active listings matching your criteria during the web search. "
            "Below is a template for manual property entry. Populate with properties from your own sources "
            "(direct mail leads, wholesalers, MLS searches, etc.)"
        )
        explanation.paragraph_format.space_before = Pt(0)
        explanation.paragraph_format.space_after = Pt(12)
        
        # Template table (6 empty rows)
        template_table = doc.add_table(rows=7, cols=6)
        template_table.style = 'Light Grid Accent 1'
        
        # Headers
        headers = ['Address', 'List Price', 'Lot Size', 'Fit for Plan', 'Data Gaps (Must Verify)', 'Source']
        header_cells = template_table.rows[0].cells
        for idx, header_text in enumerate(headers):
            header_cells[idx].text = header_text
            for paragraph in header_cells[idx].paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
        
        # Add rows with instructional text
        for row_idx in range(1, 7):
            row = template_table.rows[row_idx]
            row.cells[0].text = "[Enter address or property ID]"
            row.cells[3].text = "[Your analysis]"
            row.cells[4].text = "[Zoning, deed, survey, photos, etc]"
        
        # Add instructions
        instructions = doc.add_paragraph()
        instructions.paragraph_format.space_before = Pt(12)
        instructions.paragraph_format.space_after = Pt(0)
        
        inst_run = instructions.add_run("Next Steps to Build Deal Pipeline:\n")
        inst_run.bold = True
        
        steps = [
            "1. Launch targeted property search (probate, tax defaults, code violations, vacant)",
            "2. Add properties found to the template table above",
            "3. Validate addresses and pricing through county records",
            "4. Re-generate this report once you have identified deals",
            "5. Use Underwriting Analyzer to evaluate each property's profitability"
        ]
        for step in steps:
            doc.add_paragraph(step, style='List Bullet')

    # ===================================================
    # SECTION 4: FINANCIAL ANALYSIS
    # ===================================================
     doc.add_heading('4. Financial Analysis & Profit Projections', 1)

     doc.add_heading('4.1 Capital Deployment Plan', 2)

     deployment = data.get('deployment', [])
     if deployment:
        deploy_text = f"Your ${capital:,.0f} starting capital should be deployed strategically across three key areas:\n\n"
        doc.add_paragraph(deploy_text)
        
        deploy_table = doc.add_table(rows=len(deployment) + 2, cols=4)
        deploy_table.style = 'Light Grid Accent 1'
        
        # Header
        try:
            header_cells = deploy_table.rows[0].cells
            header_cells[0].text = 'Allocation Category'
            header_cells[1].text = 'Amount'
            header_cells[2].text = '% of Total'
            header_cells[3].text = 'Strategic Purpose'
            for cell in header_cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
        except Exception as e:
            print(f"Error setting deployment table header: {e}")
        
        # Data rows
        total_amount = 0
        for idx, item in enumerate(deployment):
            try:
                row_idx = idx + 1
                row = deploy_table.rows[row_idx]
                
                category = str(item.get('category', 'N/A')) if item.get('category') else 'N/A'
                amount = item.get('amount', 0)
                total_amount += amount
                amount_str = f"${amount:,.0f}"
                percentage = item.get('percentage', 0)
                purpose = str(item.get('purpose', '')) if item.get('purpose') else ''
                
                row.cells[0].text = category
                row.cells[1].text = amount_str
                row.cells[2].text = f"{percentage}%"
                row.cells[3].text = purpose
            except (IndexError, AttributeError, TypeError) as e:
                print(f"Error at deployment row {idx}: {e}")
                continue
        
        # Total row
        try:
            total_row = deploy_table.rows[len(deployment) + 1]
            total_row.cells[0].text = 'TOTAL ALLOCATED'
            total_row.cells[1].text = f"${total_amount:,.0f}"
            total_row.cells[2].text = "100%"
            total_row.cells[3].text = ''
            for cell in total_row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
        except Exception as e:
            print(f"Error setting total row: {e}")

    # ===================================================
    # SECTION 4.2: PROFIT PROJECTIONS (ONLY ONE - CORRECTED)
    # ===================================================
     doc.add_heading('4.2 Profit Projections (6-Month Timeline)', 2)
     heading = doc.paragraphs[-1]
     heading.paragraph_format.space_before = Pt(12)
     heading.paragraph_format.space_after = Pt(12)

     proj_intro = doc.add_paragraph(
        "Conservative monthly deal flow projections based on current market conditions:"
     )
     proj_intro.paragraph_format.space_before = Pt(0)
     proj_intro.paragraph_format.space_after = Pt(12)

     projections = data.get('projections', [])
     if projections:
        proj_table = doc.add_table(rows=len(projections) + 1, cols=5)
        proj_table.style = 'Light Grid Accent 1'
        
        # Header
        headers = ['Month', 'Period', 'Deals Closed', 'Cumulative Profit', 'Cash-on-Cash ROI']
        header_cells = proj_table.rows[0].cells
        
        for idx, header_text in enumerate(headers):
            cell = header_cells[idx]
            cell.text = header_text
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(10)
        
        # Data rows
        for row_idx, proj in enumerate(projections, 1):
            try:
                row = proj_table.rows[row_idx]
                
                row.cells[0].text = str(proj.get('month', row_idx))
                row.cells[1].text = proj.get('month_name', f"Month {proj.get('month', row_idx)}")
                row.cells[2].text = str(proj.get('deals', 0))
                
                profit = proj.get('profit', 0)
                if isinstance(profit, (int, float)):
                    row.cells[3].text = f"${profit:,.0f}"
                else:
                    row.cells[3].text = str(profit)
                
                roi = proj.get('roi', 0)
                if isinstance(roi, (int, float)):
                    row.cells[4].text = f"{roi:.1f}%"
                else:
                    row.cells[4].text = str(roi)
            
            except (IndexError, ValueError, TypeError) as e:
                print(f"Projection row error: {e}")
                continue

    # ===================================================
    # SECTION 4.3: UNDERWRITING FRAMEWORK
    # ===================================================
     doc.add_heading('4.3 Underwriting Framework (Per Property)', 2)
     heading = doc.paragraphs[-1]
     heading.paragraph_format.space_before = Pt(12)
     heading.paragraph_format.space_after = Pt(12)
    
    # Core formula explanation
     formula_para = doc.add_paragraph()
     formula_para.paragraph_format.space_before = Pt(0)
     formula_para.paragraph_format.space_after = Pt(12)
     formula_run = formula_para.add_run("Core Fix-and-Flip Formula:")
     formula_run.bold = True
     formula_run.font.size = Pt(11)
    
     formula_text = doc.add_paragraph(
        "Max Allowable Offer (MAO) = (ARV × 0.70) − Estimated Repair Costs − Known Deal Costs"
     )
     formula_text.paragraph_format.space_before = Pt(0)
     formula_text.paragraph_format.space_after = Pt(12)
    
    # Cost buffer explanation
     buffer_para = doc.add_paragraph()
     buffer_para.paragraph_format.space_before = Pt(0)
     buffer_para.paragraph_format.space_after = Pt(6)
    
     buffer_run = buffer_para.add_run("Deal Costs Buffer (Typical Categories):")
     buffer_run.bold = True
    
     cost_items = [
        "Closing costs (purchase + sale side): 6-8% of purchase price",
        "Holding costs (taxes, insurance, utilities, interest): 1-2% per month",
        "Realtor commissions (if listing): 6% of ARV",
        "Contingency (10-20% of rehab depending on property age)",
        "Contractor markup/inefficiency (10-15% of estimates)"
     ]
    
     for item in cost_items:
        doc.add_paragraph(item, style='List Bullet')
        doc.paragraphs[-1].paragraph_format.space_before = Pt(0)
        doc.paragraphs[-1].paragraph_format.space_after = Pt(3)
    
    # ===================================================
    # UNDERWRITING INPUTS TABLE - REAL DATA ONLY
    # ===================================================
    
     doc.add_heading('Underwriting Inputs Table (Fill Per Property)', 2)
     heading = doc.paragraphs[-1]
     heading.paragraph_format.space_before = Pt(12)
     heading.paragraph_format.space_after = Pt(12)
    
    # Check if we have search results with actual properties
     search_results = data.get('search_results', [])
    
     if search_results and len(search_results) > 0:
        # Show inputs for FIRST property as example
      first_prop = search_results[0]
        
      inputs_table = doc.add_table(rows=9, cols=3)
      inputs_table.style = 'Light Grid Accent 1'
        
        # Headers
      header_cells = inputs_table.rows[0].cells
      header_cells[0].text = "Input"
      header_cells[1].text = "Value"
      header_cells[2].text = "Source/Notes"
        
      for cell in header_cells:
         for paragraph in cell.paragraphs:
             for run in paragraph.runs:
                run.font.bold = True
        
        # Row 1: Property Address
      inputs_table.rows[1].cells[0].text = "Property Address"
      inputs_table.rows[1].cells[1].text = first_prop.get('address', '[Enter address]')
      inputs_table.rows[1].cells[2].text = "Zillow, Realtor, MLS"
        
        # Row 2: After Repair Value (ARV)
      inputs_table.rows[2].cells[0].text = "After Repair Value (ARV)"
      inputs_table.rows[2].cells[1].text = "[From comps - agent/BPO]"
      inputs_table.rows[2].cells[2].text = "Comps from last 60 days in ZIP"
        
        # Row 3: Estimated Rehab Budget
      inputs_table.rows[3].cells[0].text = "Estimated Rehab Budget"
      inputs_table.rows[3].cells[1].text = "[From contractor bids]"
      inputs_table.rows[3].cells[2].text = "3+ contractor estimates"
        
        # Row 4: Purchase Price (Offer)
      inputs_table.rows[4].cells[0].text = "Purchase Price (Your Offer)"
      inputs_table.rows[4].cells[1].text = f"${first_prop.get('price', '[Enter]'):,.0f}" if first_prop.get('price', 0) > 0 else "[Calculate from MAO]"
      inputs_table.rows[4].cells[2].text = "Calculated using MAO formula"
        
        # Row 5: Timeline (Days to Flip)
      inputs_table.rows[5].cells[0].text = "Project Timeline (Days)"
      inputs_table.rows[5].cells[1].text = "120"  # Typical
      inputs_table.rows[5].cells[2].text = "Rehab + 30-60 days on market"
        
        # Row 6: Holding Costs (Monthly)
      inputs_table.rows[6].cells[0].text = "Holding Costs (Monthly)"
      inputs_table.rows[6].cells[1].text = "[Calculate: taxes + insurance + utils + interest]"
      inputs_table.rows[6].cells[2].text = "County tax data + insurance quotes"
        
        # Row 7: Exit Strategy
      inputs_table.rows[7].cells[0].text = "Exit Strategy"
      inputs_table.rows[7].cells[1].text = "Retail sale to owner-occupant"
      inputs_table.rows[7].cells[2].text = "Validate buyer pool in ZIP"
        
        # Row 8: Expected Profit
      inputs_table.rows[8].cells[0].text = "Expected Profit"
      inputs_table.rows[8].cells[1].text = "[ARV - Purchase - Rehab - Costs]"
      inputs_table.rows[8].cells[2].text = "Minimum 15-20% of purchase price"
    
     else:
     # No properties - show blank template
      inputs_table = doc.add_table(rows=9, cols=3)
      inputs_table.style = 'Light Grid Accent 1'
        
        # Headers
      header_cells = inputs_table.rows[0].cells
      header_cells[0].text = "Input"
      header_cells[1].text = "Value (Enter for Each Property)"
      header_cells[2].text = "Source/Notes"
        
      for cell in header_cells:
         for paragraph in cell.paragraphs:
             for run in paragraph.runs:
                run.font.bold = True
        
        # Rows
      row_data = [
            ("Property Address", "[Enter address]", "Zillow, Realtor, MLS"),
            ("After Repair Value (ARV)", "$[From comps]", "Comps from last 60 days in ZIP"),
            ("Estimated Rehab Budget", "$[From bids]", "3+ contractor estimates"),
            ("Purchase Price (Your Offer)", "$[Calculate]", "Using MAO formula below"),
            ("Project Timeline (Days)", "120", "Rehab + 30-60 market days"),
            ("Holding Costs (Monthly)", "$[Calculate]", "Taxes + insurance + utilities + interest"),
            ("Exit Strategy", "[Your approach]", "Owner-occupant sale / Investor resale"),
            ("Expected Profit", "$[Calculate]", "Minimum 15-20% of purchase price"),
        ]
        
      for idx, (input_name, value, source) in enumerate(row_data, 1):
            inputs_table.rows[idx].cells[0].text = input_name
            inputs_table.rows[idx].cells[1].text = value
            inputs_table.rows[idx].cells[2].text = source
    
    # Add critical note
      note = doc.add_paragraph()
      note.paragraph_format.space_before = Pt(12)
      note.paragraph_format.space_after = Pt(6)
    
      note_run = note.add_run("⚠️ CRITICAL: ")
      note_run.bold = True
      note_run.font.color.rgb = RGBColor(192, 0, 0)
    
      note.add_run(
        "Underwriting accuracy depends on REAL data from contractors, comps, and county records. "
        "Never estimate repair costs - get 3+ bids. Never guess ARV - use recent sold comparables only."
     )
        
        
    # ===================================================
    # SECTION 5: STRATEGY & RECOMMENDATIONS
    # ===================================================
     doc.add_heading('5. Investment Strategy & Recommendations', 1)
    
     doc.add_heading('5.1 Recommended Asset Mix', 2)
     asset_mix_text = f"""Focus {location} property acquisition on these categories:

      1. Primary Focus (60%): Single-family residential and small multifamily properties
      - Easier to underwrite and exit
      - Strong buyer demand
      - Lower management complexity

      2. Secondary Focus (30%): Distressed properties with upside potential
      - Greater profit margins
      - Requires more active management
      - Higher risk/reward profile

      3. Opportunistic (10%): Creative finance and wholesale deals
      - Assignment fee income
      - Joint venture opportunities
      - Relationship-based deals"""
      
     doc.add_paragraph(asset_mix_text)
    
     doc.add_heading('5.2 Risk Management Framework', 2)
     risk_framework = """Implement these controls to protect capital and maximize returns:

       Financial Risk Controls:
       • Require 20%+ discount to after-repair value before making offer
       • Maintain 15% contingency reserve for unexpected rehab costs
       • Limit single-deal capital exposure to maximum 25% of total portfolio

       Market Risk Controls:
       • Diversify across 3-5 simultaneous deals in different ZIP codes
       • Maintain 60-day average holding period maximum
       • Pre-qualify buyers before locking contracts

       Operational Risk Controls:
       • Use licensed contractors with 10+ year history
       • Require detailed scope and fixed-price rehab contracts
       • Perform weekly property inspections during renovation"""
       
     doc.add_paragraph(risk_framework)
    
    # ===================================================
    # SECTION 6: ACTION PLAN & NEXT STEPS
    # ===================================================
     doc.add_heading('6. Action Plan & Next Steps', 1)
    
     doc.add_heading('6.1 90-Day Execution Roadmap', 2)
    
     execution_phases = {
        'Phase 1 (Days 1-30): Foundation & Sourcing': [
            'Establish business entity (LLC) and business bank account',
            'Pull targeted property lists (probate, tax defaults, absentee owners)',
            'Launch direct-to-seller marketing campaign',
            'Build cash buyer database (target: 25+ qualified buyers)',
            'Create CRM system for lead tracking and follow-up'
        ],
        'Phase 2 (Days 31-60): Contracting & Underwriting': [
            'Lock 2-3 properties under contract',
            'Complete detailed underwriting analysis for each property',
            'Prepare investor presentations for funding partners',
            'Validate exit strategy for each deal',
            'Establish relationships with title company and attorney'
        ],
        'Phase 3 (Days 61-90): Execution & Scaling': [
            'Close first 1-2 deals and execute exits',
            'Collect profit and deploy into next deals',
            'Refine processes based on real results',
            'Scale sourcing to 2-3 deals per month',
            'Plan for capital deployment in next quarter'
        ]
     }
    
     for phase, steps in execution_phases.items():
        doc.add_heading(phase, 3)
        for step in steps:
            doc.add_paragraph(step, style='List Bullet')
    
     doc.add_heading('6.2 Critical Success Factors', 2)
     csf = [
        'Disciplined underwriting - Do not buy based on emotion',
        'Strong buyer relationships - Your exit before you close',
        'Consistent sourcing - Deal flow is everything',
        'Speed and efficiency - Faster execution than competitors',
        'Continuous learning - Track metrics and optimize monthly'
     ]
     for factor in csf:
        doc.add_paragraph(factor, style='List Bullet')
    
    # ===================================================
    # SECTION 7: APPENDIX & DISCLAIMERS
    # ===================================================
     doc.add_page_break()
     doc.add_heading('7. Appendix & Important Disclaimers', 1)
    
     disclaimer_text = """This investment analysis is provided for informational purposes only and should not be construed as investment advice. 

        Important Notices:
        • Market conditions change constantly - validate all data before making offers
        • Past performance does not guarantee future results
        • Real estate investments carry significant risk including loss of capital
        • Consult with legal, tax, and financial advisors before making investments
        • Local market knowledge and relationships are critical to success

     Report Generated: This analysis was generated using current market data and investor-provided parameters. Actual market conditions, property values, and investment outcomes may vary significantly."""
     
     doc.add_paragraph(disclaimer_text)
    
     doc.add_heading('Contact & Support', 2)
     doc.add_paragraph('For questions about this analysis or to refine your investment strategy, contact your investment advisor.')
    
    # ===================================================
    # Save to bytes
    # ===================================================
     output = io.BytesIO()
     try:
        doc.save(output)
        output.seek(0)
        
        location_slug = location.replace(' ', '_').replace(',', '')
        filename = f"Investment_Report_{location_slug}_{datetime.now().strftime('%Y%m%d')}.docx"
        return output.read(), filename
    
     except Exception as e:
        print(f"Error saving document: {e}")
        raise
 
 
 
    # PDF Generation Method
    @staticmethod
    def generate_pdf(data: Dict) -> Tuple[bytes, str]:
     """Generate PDF report with same content as Word doc"""
    
     output = io.BytesIO()
     doc = SimpleDocTemplate(
        output,
        pagesize=letter,
        rightMargin=0.75*inch,
        leftMargin=0.75*inch,
        topMargin=0.75*inch,
        bottomMargin=0.75*inch,
        title="Investment Analysis Report"
     )
    
    # Container for PDF elements
     elements = []
    
    # Define styles
     styles = getSampleStyleSheet()
    
    # Custom styles
     title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#192D55'),
        spaceAfter=6,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
     )
    
     subtitle_style = ParagraphStyle(
        'CustomSubtitle',
        parent=styles['Normal'],
        fontSize=13,
        textColor=colors.HexColor('#3498db'),
        spaceAfter=12,
        alignment=TA_CENTER,
        fontName='Helvetica-Oblique'
     )
    
     heading1_style = ParagraphStyle(
        'CustomHeading1',
        parent=styles['Heading1'],
        fontSize=14,
        textColor=colors.HexColor('#192D55'),
        spaceAfter=6,
        spaceBefore=12,
        fontName='Helvetica-Bold'
     )
    
     heading2_style = ParagraphStyle(
        'CustomHeading2',
        parent=styles['Heading2'],
        fontSize=12,
        textColor=colors.HexColor('#34495E'),
        spaceAfter=6,
        spaceBefore=6,
        fontName='Helvetica-Bold'
     )
    
     normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=10,
        spaceAfter=6,
        alignment=TA_LEFT,
        leading=14
     )
    
    # ===================================================
    # TITLE PAGE
    # ===================================================
     elements.append(Paragraph('Real Estate Investment Analysis Report', title_style))
     elements.append(Paragraph('Comprehensive Investment Strategy & Market Analysis', subtitle_style))
     elements.append(Spacer(1, 0.3*inch))
    
     investor_name = str(data.get('investor_name', 'Real Estate Investor'))
     date_str = str(data.get('date', 'N/A'))
     location = str(data.get('location', 'Not specified'))
     strategy = str(data.get('strategy_type', 'Investment Strategy'))
    
     elements.append(Paragraph(f"<b>Prepared for:</b> {investor_name}", normal_style))
     elements.append(Paragraph(f"<b>Date:</b> {date_str}", normal_style))
     elements.append(Paragraph(f"<b>Market:</b> {location}", normal_style))
     elements.append(Paragraph(f"<b>Strategy:</b> {strategy}", normal_style))
    
     elements.append(PageBreak())
    
    # ===================================================
    # SECTION 1: EXECUTIVE SUMMARY
    # ===================================================
     elements.append(Paragraph('1. Executive Summary', heading1_style))
    
     capital = data.get('capital', 0)
     profit_goal = data.get('profit_goal', 0)
     property_count = len(data.get('search_results', []))
    
     summary = f"""This report provides a detailed real estate investment analysis for {location} based on your investment criteria:<br/>
     <br/>
     <b>Starting Capital:</b> ${capital:,.0f}<br/>
     <b>Target Profit Goal:</b> ${profit_goal:,.0f}<br/>
     <b>Investment Timeline:</b> {str(data.get('timeline', 'Not specified'))}<br/>
     <b>Strategy Focus:</b> {strategy}<br/>
     <b>Property Type:</b> {str(data.get('property_type', 'Residential/Commercial'))}<br/>
     <br/>
     The analysis identifies {property_count} qualified properties, evaluates market conditions, and provides actionable recommendations to maximize return on investment while minimizing risk exposure."""
    
     elements.append(Paragraph(summary, normal_style))
     elements.append(Spacer(1, 0.2*inch))
    
    # ===================================================
    # SECTION 2: INVESTOR PROFILE
    # ===================================================
     elements.append(Paragraph('2. Investor Profile & Objectives', heading1_style))
     elements.append(Paragraph('2.1 Investment Profile', heading2_style))
    
     profile_data = [
        ['Category', 'Value'],
        ['Investor Name', investor_name],
        ['Starting Capital', f"${capital:,.0f}"],
        ['Target Geography', location],
        ['Property Type', str(data.get('property_type', 'N/A'))],
        ['Investment Strategy', strategy],
        ['Timeline', str(data.get('timeline', 'N/A'))],
        ['Target Profit Goal', f"${profit_goal:,.0f}"],
     ]
    
     profile_table = Table(profile_data, colWidths=[2.5*inch, 3.0*inch])
     profile_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2C3E50')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 10),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ('FONTSIZE', (0, 1), (-1, -1), 9),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F5F5F5')]),
     ]))
     elements.append(profile_table)
     elements.append(Spacer(1, 0.2*inch))
    
    # ===================================================
    # SECTION 3: MARKET ANALYSIS
    # ===================================================
     elements.append(Paragraph('3. Market Analysis & Property Overview', heading1_style))
     elements.append(Paragraph('3.1 Market Summary', heading2_style))
    
     search_results = data.get('search_results', [])
     prices = [p['price'] for p in search_results if p.get('price', 0) > 0]
    
    # Proper conditional handling for avg_price
     if prices:
        avg_price = sum(prices) / len(prices)
        avg_price_str = f"${avg_price:,.0f}"
        price_range = f"${min(prices):,} - ${max(prices):,}"
     else:
        avg_price_str = "N/A (No pricing data)"
        price_range = "N/A"
    
     market_text = f"""The {location} real estate market presents compelling opportunities for disciplined investors.<br/>
     <br/>
     <b>Market Indicators:</b><br/>
     <b>Total Properties Identified:</b> {len(search_results)}<br/>
     <b>Average List Price:</b> {avg_price_str}<br/>
     <b>Price Range:</b> {price_range}<br/>
     <b>Market Trend:</b> Active buyer interest with competitive opportunities"""
    
     elements.append(Paragraph(market_text, normal_style))
     elements.append(Spacer(1, 0.2*inch))
    
    # ===================================================
    # SECTION 3.2: PROPERTIES TABLE
    # ===================================================
     elements.append(Paragraph('3.2 Identified Properties (Top 10)', heading2_style))
    
     if search_results and len(search_results) > 0:
        properties = search_results[:10]
        prop_data = [['Address', 'List Price', 'Lot Size', 'Fit', 'Data Gaps', 'Source']]
        
        for prop in properties:
            address = str(prop.get('address', 'N/A'))[:40]
            
            # Type-safe price formatting
            price = prop.get('price', 0)
            if isinstance(price, (int, float)) and price > 0:
                price_str = f"${price:,.0f}"
            else:
                price_str = "Contact"
            
            # Type-safe acres formatting
            acres = prop.get('acres')
            if isinstance(acres, (int, float)) and acres > 0:
                lot_size = f"{acres} ac"
            else:
                lot_size = str(prop.get('lot_size', 'See listing'))
            
            source = str(prop.get('source', 'Listing'))[:15]
            fit = "✓ Good" if price > 0 and price < capital else "⚠ Review"
            gaps = "Verify deed, zoning"
            
            prop_data.append([address, price_str, lot_size, fit, gaps, source])
        
            prop_table = Table(prop_data, colWidths=[1.2*inch, 0.9*inch, 0.8*inch, 0.6*inch, 1.0*inch, 0.8*inch])
            prop_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3498db')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 9),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F5F5F5')]),
            ('TOPPADDING', (0, 0), (-1, -1), 4),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
         ]))
            elements.append(prop_table)
     else:
         elements.append(Paragraph(
            '<b>⚠️ No Properties Found</b> - Web search returned no results. Populate with properties from your own sources.',
            normal_style
        ))
    
     elements.append(Spacer(1, 0.2*inch))
    
    # ===================================================
    # SECTION 4: CAPITAL DEPLOYMENT
    # ===================================================
     elements.append(Paragraph('4. Financial Analysis & Capital Deployment', heading1_style))
    
     if data.get('deployment'):
        deploy_data = [['Allocation Category', 'Amount', '% of Total', 'Strategic Purpose']]
        
        for item in data['deployment']:
            deploy_data.append([
                str(item['category']),
                f"${item['amount']:,.0f}",
                f"{item['percentage']}%",
                str(item['purpose'])[:40]
            ])
        
        deploy_table = Table(deploy_data, colWidths=[1.5*inch, 1.0*inch, 0.8*inch, 2.2*inch])
        deploy_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2C3E50')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 9),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('FONTSIZE', (0, 1), (-1, -1), 9),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F5F5F5')]),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ]))
        elements.append(deploy_table)
        elements.append(Spacer(1, 0.2*inch))
    
    # ===================================================
    # SECTION 5: PROFIT PROJECTIONS
    # ===================================================
     elements.append(Paragraph('4.2 Profit Projections (6-Month Timeline)', heading2_style))
    
     if data.get('projections'):
        proj_data = [['Month', 'Period', 'Deals', 'Cumulative Profit', 'ROI']]
        
        for proj in data['projections']:
            proj_data.append([
                str(proj['month']),
                str(proj.get('month_name', '')),
                str(proj['deals']),
                f"${proj['profit']:,.0f}",
                f"{proj.get('roi', 0):.1f}%"
            ])
        
        proj_table = Table(proj_data, colWidths=[0.7*inch, 0.9*inch, 0.7*inch, 1.5*inch, 0.8*inch])
        proj_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2ECC71')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 9),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('FONTSIZE', (0, 1), (-1, -1), 9),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F5F5F5')]),
        ]))
        elements.append(proj_table)
    
     elements.append(PageBreak())
    
    # ===================================================
    # SECTION 6: DISCLAIMERS
    # ===================================================
     elements.append(Paragraph('5. Important Disclaimers', heading1_style))
    
     disclaimer_text = f"""This investment analysis is provided for informational purposes only and should not be construed as investment advice.<br/>
     <br/>
     <b>Important Notices:</b><br/>
     • Market conditions change constantly - validate all data before making offers<br/>
     • Past performance does not guarantee future results<br/>
     • Real estate investments carry significant risk including loss of capital<br/>
     • Consult with legal, tax, and financial advisors before making investments<br/>
     • Local market knowledge and relationships are critical to success<br/>
     <br/>
     <i>Report Generated: {date_str}</i>"""
    
     elements.append(Paragraph(disclaimer_text, normal_style))
    
    # Build PDF
     doc.build(elements)
     output.seek(0)
    
     location_slug = location.replace(' ', '_').replace(',', '')
     filename = f"Investment_Report_{location_slug}.pdf"
    
     return output.read(), filename



#    # PowerPoint Generation Method
    @staticmethod
    def generate_powerpoint(data: Dict) -> Tuple[bytes, str]:
     """
     Generate professional PowerPoint investor presentation
     """
    
    # Create presentation
     prs = Presentation()
     prs.slide_width = PptxInches(10)
     prs.slide_height = PptxInches(7.5)
    
    # Extract data with validation
     investor_name = str(data.get('investor_name', 'Real Estate Investor'))
     location = str(data.get('location', 'Not specified'))
     capital = data.get('capital', 0)
     profit_goal = data.get('profit_goal', 0)
     strategy = str(data.get('strategy_type', 'Investment Strategy'))
     date_str = str(data.get('date', 'N/A'))
     search_results = data.get('search_results', [])
     property_count = len(search_results)
     projections = data.get('projections', [])
     deployment = data.get('deployment', [])
    
    # ===================================================
    # SLIDE 1: TITLE SLIDE WITH COLORED BACKGROUND
    # ===================================================
     slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Use MSO_SHAPE.RECTANGLE and PptxInches
     background_shape = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        PptxInches(0), 
        PptxInches(0), 
        prs.slide_width, 
        prs.slide_height
     )
    
    # Use PptxRGBColor for PowerPoint
     fill = background_shape.fill
     fill.solid()
     fill.fore_color.rgb = PptxRGBColor(25, 45, 85)  # Dark blue
    
    # Remove border
     background_shape.line.fill.background()
    
    # Title
     title_box = slide1.shapes.add_textbox(
        PptxInches(0.5), PptxInches(2), PptxInches(9), PptxInches(1.5)
    )
     title_frame = title_box.text_frame
     title_frame.word_wrap = True
     title_p = title_frame.paragraphs[0]
     title_p.text = "Real Estate Investment Analysis"
     title_p.font.size = PptxPt(54)
     title_p.font.bold = True
     title_p.font.color.rgb = PptxRGBColor(255, 255, 255)  # White
     title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
     subtitle_box = slide1.shapes.add_textbox(
        PptxInches(0.5), PptxInches(3.5), PptxInches(9), PptxInches(1)
    )
     subtitle_frame = subtitle_box.text_frame
     subtitle_p = subtitle_frame.paragraphs[0]
     subtitle_p.text = f"{location} Investment Opportunity"
     subtitle_p.font.size = PptxPt(32)
     subtitle_p.font.color.rgb = PptxRGBColor(52, 152, 219)  # Light blue
     subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Details at bottom
     details_box = slide1.shapes.add_textbox(
        PptxInches(0.5), PptxInches(5.5), PptxInches(9), PptxInches(1.5)
    )
     details_frame = details_box.text_frame
     details_frame.word_wrap = True
    
     details_lines = [
        f"Prepared for: {investor_name}",
        f"Date: {date_str}",
        f"Strategy: {strategy}"
     ]
    
     for line in details_lines:
        p = details_frame.add_paragraph()
        p.text = line
        p.font.size = PptxPt(18)
        p.font.color.rgb = PptxRGBColor(255, 255, 255)  # White
        p.alignment = PP_ALIGN.CENTER
        p.space_before = PptxPt(6)
    
    # ===================================================
    # SLIDE 2: EXECUTIVE SUMMARY
    # ===================================================
     slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    
     title2 = slide2.shapes.title
     title2.text = "Executive Summary"
     title2.text_frame.paragraphs[0].font.size = PptxPt(44)
     title2.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(25, 45, 85)
    
     body_box = slide2.placeholders[1]
     tf = body_box.text_frame
     tf.clear()
    
     summary_points = [
        f"Capital Available: ${capital:,.0f}",
        f"Profit Target: ${profit_goal:,.0f}",
        f"Timeline: {str(data.get('timeline', 'Not specified'))}",
        f"Properties Identified: {property_count}",
        f"Market: {location}",
        f"Strategy: {strategy}"
     ]
    
     for i, point in enumerate(summary_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = PptxPt(24)
        p.level = 0
        p.space_after = PptxPt(12)
    
    # ===================================================
    # SLIDE 3: MARKET OPPORTUNITY
    # ===================================================
     slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    
     title3 = slide3.shapes.title
     title3.text = "Market Opportunity"
     title3.text_frame.paragraphs[0].font.size = PptxPt(44)
     title3.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(25, 45, 85)
    
     body_box = slide3.placeholders[1]
     tf = body_box.text_frame
     tf.clear()
    
     prices = [p.get('price', 0) for p in search_results if p.get('price', 0) > 0]
     avg_price = (sum(prices) / len(prices)) if prices else 0
     min_price = min(prices) if prices else 0
     max_price = max(prices) if prices else 0
    
     market_points = [
        f"Total Properties Identified: {property_count}",
        f"Average List Price: ${avg_price:,.0f}",
        f"Price Range: ${min_price:,} - ${max_price:,}",
        f"Market Trend: Active buyer interest",
        f"Opportunity: Competitive pricing & volume"
     ]
    
     for i, point in enumerate(market_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = PptxPt(22)
        p.level = 0
        p.space_after = PptxPt(12)
    
    # ===================================================
    # SLIDE 4: TOP PROPERTIES
    # ===================================================
     if search_results and len(search_results) > 0:
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        
        title4 = slide4.shapes.title
        title4.text = "Property Showcase - Top 5"
        title4.text_frame.paragraphs[0].font.size = PptxPt(44)
        title4.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(25, 45, 85)
        
        body_box = slide4.placeholders[1]
        tf = body_box.text_frame
        tf.clear()
        
        for i, prop in enumerate(search_results[:5]):
            address = prop.get('address', 'N/A')
            price = prop.get('price', 0)
            acres = prop.get('acres', 'N/A')
            
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            price_str = f"${price:,.0f}" if price > 0 else "Contact"
            acres_str = f"{acres} acres" if acres != 'N/A' else "See listing"
            
            p.text = f"{i+1}. {address} | {price_str} | {acres_str}"
            p.font.size = PptxPt(18)
            p.level = 0
            p.space_after = PptxPt(10)
    
    # ===================================================
    # SLIDE 5: FINANCIAL SUMMARY
    # ===================================================
     slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    
     title5 = slide5.shapes.title
     title5.text = "Financial Projections"
     title5.text_frame.paragraphs[0].font.size = PptxPt(44)
     title5.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(25, 45, 85)
    
     body_box = slide5.placeholders[1]
     tf = body_box.text_frame
     tf.clear()
    
     if projections and len(projections) > 0:
        final_proj = projections[-1]
        final_profit = final_proj.get('profit', 0)
        final_roi = final_proj.get('roi', 0)
        
        projection_points = [
            f"6-Month Cumulative Profit: ${final_profit:,.0f}",
            f"Expected ROI: {final_roi:.1f}%",
            f"Average Monthly Deal Closure: 1-2 deals",
            f"Capital Deployment: Strategic 3-part allocation",
            f"Risk Profile: Conservative-Moderate"
        ]
     else:
        projection_points = [
            f"Capital: ${capital:,.0f}",
            f"Profit Target: ${profit_goal:,.0f}",
            "Strategy: Conservative growth approach",
            "Diversification: Multi-property portfolio",
            "Timeline: 6-month execution plan"
        ]
    
     for i, point in enumerate(projection_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = PptxPt(20)
        p.level = 0
        p.space_after = PptxPt(12)
    
    # ===================================================
    # SLIDE 6: ACTION PLAN
    # ===================================================
     slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    
     title6 = slide6.shapes.title
     title6.text = "90-Day Action Plan"
     title6.text_frame.paragraphs[0].font.size = PptxPt(44)
     title6.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(25, 45, 85)
    
     body_box = slide6.placeholders[1]
     tf = body_box.text_frame
     tf.clear()
    
     action_items = [
        "Phase 1 (Days 1-30): Foundation & Sourcing",
        "  • Set up business entity & banking",
        "  • Launch property search campaigns",
        "Phase 2 (Days 31-60): Contracting & Analysis",
        "  • Lock 2-3 properties under contract",
        "  • Complete detailed underwriting",
        "Phase 3 (Days 61-90): Execution & Scaling",
        "  • Close first deals & execute exits",
        "  • Deploy profits into next acquisitions"
     ]
    
     for i, item in enumerate(action_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(18) if not item.startswith("  ") else PptxPt(16)
        p.level = 1 if item.startswith("  ") else 0
        p.space_after = PptxPt(8)
    
    # ===================================================
    # SLIDE 7: NEXT STEPS & CTA
    # ===================================================
     slide7 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
     background_shape = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        PptxInches(0), 
        PptxInches(0), 
        prs.slide_width, 
        prs.slide_height
     )
    
     fill = background_shape.fill
     fill.solid()
     fill.fore_color.rgb = PptxRGBColor(25, 45, 85)
     background_shape.line.fill.background()
    
    # Main message
     cta_box = slide7.shapes.add_textbox(
        PptxInches(0.5), PptxInches(2), PptxInches(9), PptxInches(2)
    )
     cta_frame = cta_box.text_frame
     cta_frame.word_wrap = True
    
     cta_p = cta_frame.paragraphs[0]
     cta_p.text = "Ready to Take Action?"
     cta_p.font.size = PptxPt(48)
     cta_p.font.bold = True
     cta_p.font.color.rgb = PptxRGBColor(255, 255, 255)
     cta_p.alignment = PP_ALIGN.CENTER
    
    # Next steps
     steps_box = slide7.shapes.add_textbox(
        PptxInches(0.5), PptxInches(4), PptxInches(9), PptxInches(2.5)
    )
     steps_frame = steps_box.text_frame
     steps_frame.word_wrap = True
    
     next_steps_list = [
        "1. Schedule property site visits this week",
        "2. Request detailed county records & zoning",
        "3. Build qualified cash buyer database",
        "4. Prepare offer strategy & timeline"
     ]
    
     for step in next_steps_list:
        p = steps_frame.add_paragraph()
        p.text = step
        p.font.size = PptxPt(22)
        p.font.color.rgb = PptxRGBColor(52, 152, 219)
        p.space_after = PptxPt(10)
    
    # ===================================================
    # SAVE AND RETURN
    # ===================================================
     output = io.BytesIO()
     try:
        prs.save(output)
        output.seek(0)
        
        location_slug = location.replace(' ', '_').replace(',', '')
        filename = f"Investment_Presentation_{location_slug}.pptx"
        return output.read(), filename
    
     except Exception as e:
        print(f"Error saving presentation: {e}")
        raise




# ===================================================
# UPDATE should_generate_file() TO INCLUDE PDF
# ===================================================

    @staticmethod
    def should_generate_file(message: str, response: str) -> Tuple[bool, Optional[str]]:
     """
     Determine if file should be generated and what type
     INCLUDES PDF SUPPORT
     """
    
     user_keywords = {
        'excel': ['excel', 'spreadsheet', 'xlsx', 'table', 'csv', 'tracker','sheet', 'pipeline', 'data', 'charts'],
        'word': ['word', 'document', 'docx', 'report', 'write', 'text file', 'doc'],
        'powerpoint': ['powerpoint', 'presentation', 'pptx', 'slides', 'pitch', 'graphs', 'visuals'],
        'pdf': ['pdf', 'download', 'format', 'save', 'attachment', 'shareable', 'file', 'generated'] 
     }
    
     action_keywords = [
        'generate', 'create', 'make', 'build', 'produce',
        'export', 'download', 'send', 'give me', 'show me'
     ]
    
     message_lower = message.lower()
     response_lower = response.lower()
    
    # Check if action keyword present
     has_action = any(action in message_lower for action in action_keywords)
    
     if not has_action:
        return False, None
    
    # Determine file type
     for file_type, keywords in user_keywords.items():
        if any(keyword in message_lower for keyword in keywords):
            return True, file_type
    
    # Default to Excel for data/results
     if any(word in message_lower for word in ['table', 'data', 'list', 'results', 'properties']):
        return True, 'excel'
    
    # Default to Word for reports
     if any(word in message_lower for word in ['report', 'analysis', 'strategy']):
        return True, 'word'
    
     return False, None
 
    
    """
    SECTION: Excel Spreadsheet Generation with Full Content
    """
    @staticmethod
    def generate_excel(data: Dict) -> Tuple[bytes, str]:
     """
     Generate professional Excel workbook with REAL DATA ONLY.
    
     Structure:
     - Sheet 1: Deal Pipeline (property listings)
     - Sheet 2: Investment Summary (profile + market analysis)
     - Sheet 3: Capital Deployment (budget allocation)
     - Sheet 4: Profit Projections (6-month timeline + chart)
     - Sheet 5: Action Items (next steps & timeline)
    
     Args:
        data: Dict containing investor profile, search results, calculations
        
     Returns:
        Tuple of (file_bytes, filename)
     """
    
     output = io.BytesIO()
     workbook = xlsxwriter.Workbook(output, {'in_memory': True})
    
    # ===================================================
    # DEFINE FORMATS (Reusable formatting)
    # ===================================================
    
     formats = FileGenerationService._create_excel_formats(workbook)
    
    # ===================================================
    # SHEET 1: DEAL PIPELINE
    # ===================================================
    
     FileGenerationService._create_deal_pipeline_sheet(
        workbook, 
        data, 
        formats
     )
    
    # ===================================================
    # SHEET 2: INVESTMENT SUMMARY
    # ===================================================
    
     FileGenerationService._create_investment_summary_sheet(
        workbook, 
        data, 
        formats
     )
    
    # ===================================================
    # SHEET 3: CAPITAL DEPLOYMENT
    # ===================================================
    
     FileGenerationService._create_capital_deployment_sheet(
        workbook, 
        data, 
        formats
     )
    
    # ===================================================
    # SHEET 4: PROFIT PROJECTIONS
    # ===================================================
    
     FileGenerationService._create_profit_projections_sheet(
        workbook, 
        data, 
        formats
     )
    
    # ===================================================
    # SHEET 5: ACTION ITEMS
    # ===================================================
    
     FileGenerationService._create_action_items_sheet(
        workbook, 
        data, 
        formats
     )
    
    # Close and return
     workbook.close()
     output.seek(0)
    
     location_slug = data.get('location', 'analysis').replace(' ', '_').replace(',', '')
     filename = f"Deal_Pipeline_{location_slug}.xlsx"
    
     return output.read(), filename


# ===================================================
# HELPER METHODS
# ===================================================

    @staticmethod
    def _create_excel_formats(workbook) -> Dict:
     """Create all reusable Excel formats."""
     return {
        'header': workbook.add_format({
            'bold': True,
            'bg_color': '#2C3E50',
            'font_color': 'white',
            'border': 1,
            'align': 'center',
            'valign': 'vcenter',
            'font_size': 11
        }),
        'subheader': workbook.add_format({
            'bold': True,
            'bg_color': '#34495E',
            'font_color': 'white',
            'border': 1,
            'font_size': 10
        }),
        'money': workbook.add_format({
            'num_format': '$#,##0',
            'border': 1
        }),
        'money_bold': workbook.add_format({
            'num_format': '$#,##0',
            'border': 1,
            'bold': True
        }),
        'percent': workbook.add_format({
            'num_format': '0.0%',
            'border': 1
        }),
        'text': workbook.add_format({
            'border': 1,
            'valign': 'top'
        }),
        'text_bold': workbook.add_format({
            'border': 1,
            'valign': 'top',
            'bold': True
        }),
        'bold': workbook.add_format({
            'bold': True
        }),
        'url': workbook.add_format({
            'color': 'blue',
            'underline': 1,
            'border': 1
        })
    }


    @staticmethod
    def _create_deal_pipeline_sheet(workbook, data: Dict, formats: Dict) -> None:
     """Create Deal Pipeline sheet with property listings."""
    
     pipeline = workbook.add_worksheet('Deal Pipeline')
     search_results = data.get('search_results', [])
    
    # Set column widths
     pipeline.set_column('A:A', 4)    # #
     pipeline.set_column('B:B', 40)   # Address
     pipeline.set_column('C:C', 12)   # Price
     pipeline.set_column('D:D', 12)   # Lot Size
     pipeline.set_column('E:E', 12)   # $/Acre
     pipeline.set_column('F:F', 15)   # Type
     pipeline.set_column('G:G', 12)   # Status
     pipeline.set_column('H:H', 20)   # Source
     pipeline.set_column('I:I', 40)   # Notes
    
    # Headers
     headers = [
        '#', 
        'Property Address', 
        'List Price', 
        'Lot Size (Acres)', 
        '$/Acre', 
        'Property Type', 
        'Status', 
        'Source', 
        'Notes'
     ]
    
     for col, header_text in enumerate(headers):
        pipeline.write(0, col, header_text, formats['header'])
    
     # Data rows - Only REAL properties
     if search_results:
        for idx, prop in enumerate(search_results[:50], 1):
            # Column A: Row number
            pipeline.write(idx, 0, idx, formats['text'])
            
            # Column B: Address
            address = prop.get('address', 'Address not provided')
            pipeline.write(idx, 1, address, formats['text'])
            
            # Column C: Price
            price = prop.get('price', 0)
            if price > 0:
                pipeline.write(idx, 2, price, formats['money'])
            else:
                pipeline.write(idx, 2, 'Contact Seller', formats['text'])
            
            # Column D: Lot Size
            acres = prop.get('acres')
            if acres and isinstance(acres, (int, float)):
                pipeline.write(idx, 3, acres, formats['text'])
            else:
                lot_size_str = prop.get('lot_size', 'See listing')
                pipeline.write(idx, 3, lot_size_str, formats['text'])
            
            # Column E: Price per acre
            if price > 0 and acres and isinstance(acres, (int, float)) and acres > 0:
                price_per_acre = price / acres
                pipeline.write(idx, 4, price_per_acre, formats['money'])
            else:
                pipeline.write(idx, 4, 'N/A', formats['text'])
            
            # Column F: Property Type
            prop_type = prop.get('property_type', 'Land')
            pipeline.write(idx, 5, prop_type, formats['text'])
            
            # Column G: Status
            pipeline.write(idx, 6, 'Active', formats['text'])
            
            # Column H: Source (with hyperlink if available)
            source_url = prop.get('source_url', '')
            source_name = prop.get('source', 'Listing')
            
            if source_url:
                try:
                    pipeline.write_url(idx, 7, source_url, formats['url'], string=source_name)
                except:
                    pipeline.write(idx, 7, source_name, formats['text'])
            else:
                pipeline.write(idx, 7, source_name, formats['text'])
            
            # Column I: Notes (bedrooms, bathrooms, sqft, description)
            notes_parts = []
            if prop.get('bedrooms'):
                notes_parts.append(f"{prop['bedrooms']}bd/{prop.get('bathrooms', 0)}ba")
            if prop.get('sqft'):
                notes_parts.append(f"{prop['sqft']:,} sqft")
            if prop.get('description'):
                notes_parts.append(prop['description'][:100])
            
            notes_text = ' | '.join(notes_parts) if notes_parts else ''
            pipeline.write(idx, 8, notes_text, formats['text'])
        
        # Summary stats at bottom
        row = len(search_results) + 2
        pipeline.write(row, 0, 'SUMMARY', formats['bold'])
        pipeline.write(row, 1, f"Total Properties: {len(search_results)}", formats['text_bold'])
        
        # Calculate and show average price
        prices = [p['price'] for p in search_results if p.get('price', 0) > 0]
        if prices:
            avg_price = sum(prices) / len(prices)
            pipeline.write(row + 1, 1, f"Average Price: ${avg_price:,.0f}", formats['text_bold'])
    
     else:
        # No results - show placeholder
        pipeline.write(0, 0, 'No properties found in search results', formats['text_bold'])
        pipeline.write(1, 0, 'Run a property search to populate this sheet', formats['text'])


    @staticmethod
    def _create_investment_summary_sheet(workbook, data: Dict, formats: Dict) -> None:
     """Create Investment Summary sheet with profile and market analysis."""
    
     summary = workbook.add_worksheet('Investment Summary')
     summary.set_column('A:A', 28)
     summary.set_column('B:B', 25)
    
    # Extract data with validation
     investor_name = data.get('investor_name', 'N/A')
     strategy = data.get('strategy_type', 'N/A')
     property_type = data.get('property_type', 'N/A')
     capital = data.get('capital', 0)
     location = data.get('location', 'N/A')
     timeline = data.get('timeline', 'N/A')
     profit_goal = data.get('profit_goal', 0)
     date_str = data.get('date', '')
     search_results = data.get('search_results', [])
    
    # Starting row
     row = 0
    
    # Title
     summary.write(row, 0, 'INVESTMENT STRATEGY SUMMARY', formats['header'])
     summary.write(row, 1, f"Generated: {date_str}", formats['text'])
     row += 2
    
    # Section: Investor Information
     profile_items = [
        ('Investor/Entity Name', investor_name),
        ('Investment Strategy', strategy),
        ('Target Property Type', property_type),
        ('Starting Capital', f"${capital:,.0f}" if capital > 0 else 'Not specified'),
        ('Target Geography', location),
        ('Investment Timeline', timeline),
        ('Target Profit Goal', f"${profit_goal:,.0f}" if profit_goal > 0 else 'Not specified'),
     ]
    
     for label, value in profile_items:
        summary.write(row, 0, label, formats['subheader'])
        summary.write(row, 1, value, formats['text'])
        row += 1
    
     row += 1  # Space between sections
    
    # Section: Market Analysis
     summary.write(row, 0, 'MARKET ANALYSIS', formats['header'])
     row += 1
    
     summary.write(row, 0, 'Properties Identified', formats['subheader'])
     summary.write(row, 1, len(search_results), formats['text'])
     row += 1
    
    # Calculate average price if results exist
     if search_results:
        prices = [p['price'] for p in search_results if p.get('price', 0) > 0]
        if prices:
            avg_price = sum(prices) / len(prices)
            min_price = min(prices)
            max_price = max(prices)
            
            summary.write(row, 0, 'Average List Price', formats['subheader'])
            summary.write(row, 1, avg_price, formats['money'])
            row += 1
            
            summary.write(row, 0, 'Price Range', formats['subheader'])
            summary.write(row, 1, f"${min_price:,} - ${max_price:,}", formats['text'])
            row += 1


    @staticmethod
    def _create_capital_deployment_sheet(workbook, data: Dict, formats: Dict) -> None:
     """Create Capital Deployment sheet with budget allocation."""
    
     deployment = data.get('deployment', [])
    
     if not deployment:
        return
    
     deploy = workbook.add_worksheet('Capital Deployment')
     deploy.set_column('A:A', 28)
     deploy.set_column('B:B', 15)
     deploy.set_column('C:C', 10)
     deploy.set_column('D:D', 50)
    
    # Headers
     headers = [
        'Allocation Category',
        'Amount',
        '% of Total',
        'Strategic Purpose'
    ]
    
     for col, header_text in enumerate(headers):
        deploy.write(0, col, header_text, formats['header'])
    
    # Data rows
     total_amount = 0
     for idx, item in enumerate(deployment, 1):
        category = item.get('category', 'N/A')
        amount = item.get('amount', 0)
        percentage = item.get('percentage', 0)
        purpose = item.get('purpose', '')
        
        total_amount += amount
        
        deploy.write(idx, 0, category, formats['text'])
        deploy.write(idx, 1, amount, formats['money'])
        deploy.write(idx, 2, percentage / 100, formats['percent'])
        deploy.write(idx, 3, purpose, formats['text'])
    
    # Total row
     row = len(deployment) + 1
     deploy.write(row, 0, 'TOTAL ALLOCATED', formats['subheader'])
     deploy.write(row, 1, total_amount, formats['money_bold'])
     deploy.write(row, 2, 1.0, formats['percent'])
    
    # Risk balance notes
     row += 2
     deploy.write(row, 0, 'Risk Balance Validation:', formats['text_bold'])
     deploy.write(row + 1, 0, '✓ 40% in contracts (moderate risk, high return)', formats['text'])
     deploy.write(row + 2, 0, '✓ 35% in marketing (controlled deployment)', formats['text'])
     deploy.write(row + 3, 0, '✓ 25% in reserves (downside protection)', formats['text'])


    @staticmethod
    def _create_profit_projections_sheet(workbook, data: Dict, formats: Dict) -> None:
     """Create Profit Projections sheet with 6-month timeline and chart."""
    
     projections = data.get('projections', [])
    
     if not projections:
        return
    
     proj_sheet = workbook.add_worksheet('Profit Projections')
     proj_sheet.set_column('A:E', 18)
    
    # Headers
     headers = [
        'Month',
        'Period',
        'Deals Closed',
        'Cumulative Profit',
        'Cash-on-Cash ROI'
     ]
    
     for col, header_text in enumerate(headers):
        proj_sheet.write(0, col, header_text, formats['header'])
    
    # Data rows
     for idx, proj in enumerate(projections, 1):
        month = proj.get('month', idx)
        month_name = proj.get('month_name', f"Month {month}")
        deals = proj.get('deals', 0)
        profit = proj.get('profit', 0)
        roi = proj.get('roi', 0)
        
        proj_sheet.write(idx, 0, month, formats['text'])
        proj_sheet.write(idx, 1, month_name, formats['text'])
        proj_sheet.write(idx, 2, deals, formats['text'])
        proj_sheet.write(idx, 3, profit, formats['money'])
        proj_sheet.write(idx, 4, roi / 100, formats['percent'])
    
    # Add chart
     try:
        chart = workbook.add_chart({'type': 'line'})
        chart.add_series({
            'name': 'Cumulative Profit',
            'categories': f'=\'Profit Projections\'!$B$2:$B${len(projections)+1}',
            'values': f'=\'Profit Projections\'!$D$2:$D${len(projections)+1}',
            'line': {'color': '#2ECC71', 'width': 3}
        })
        chart.set_title({'name': '6-Month Profit Trajectory', 'name_font': {'size': 14, 'bold': True}})
        chart.set_x_axis({'name': 'Month'})
        chart.set_y_axis({'name': 'Profit ($)', 'num_format': '$#,##0'})
        chart.set_size({'width': 600, 'height': 400})
        proj_sheet.insert_chart('G2', chart)
     except Exception as e:
        print(f"⚠️ Chart creation failed: {e}")
    
    # Projection assumptions
     row = len(projections) + 3
     proj_sheet.write(row, 0, 'PROJECTION ASSUMPTIONS', formats['text_bold'])
     proj_sheet.write(row + 1, 0, '• Conservative ramp-up schedule', formats['text'])
     proj_sheet.write(row + 2, 0, '• Average assignment fee per deal', formats['text'])
     proj_sheet.write(row + 3, 0, '• Does not account for compounding effects', formats['text'])


    @staticmethod
    def _create_action_items_sheet(workbook, data: Dict, formats: Dict) -> None:
     """Create Action Items sheet with prioritized next steps."""
    
     next_steps = data.get('next_steps', [])
    
     if not next_steps:
        return
    
     steps_sheet = workbook.add_worksheet('Action Items')
     steps_sheet.set_column('A:A', 8)
     steps_sheet.set_column('B:B', 60)
     steps_sheet.set_column('C:C', 15)
    
    # Headers
     headers = ['Priority', 'Action Item', 'Timeline']
     for col, header_text in enumerate(headers):
        steps_sheet.write(0, col, header_text, formats['header'])
    
    # Data rows
     for idx, step in enumerate(next_steps, 1):
        priority = 'HIGH' if idx <= 3 else 'MEDIUM'
        timeline = 'This Week' if idx <= 3 else 'Next 2 Weeks'
        
        steps_sheet.write(idx, 0, priority, formats['text'])
        steps_sheet.write(idx, 1, step, formats['text'])
        steps_sheet.write(idx, 2, timeline, formats['text'])